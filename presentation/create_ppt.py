
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD    = RGBColor(0x1E, 0x29, 0x3B)
ACCENT     = RGBColor(0x38, 0xBD, 0xF8)
ACCENT2    = RGBColor(0x22, 0xD3, 0xEE)
GREEN      = RGBColor(0x4A, 0xDE, 0x80)
ORANGE     = RGBColor(0xFB, 0x92, 0x3C)
PURPLE     = RGBColor(0xA7, 0x8B, 0xFA)
RED        = RGBColor(0xF8, 0x71, 0x71)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY   = RGBColor(0x94, 0xA3, 0xB8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "  " + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    return txBox

def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    shape = add_shape_bg(slide, left, top, width, height, RGBColor(0x0D, 0x11, 0x17), RGBColor(0x30, 0x3A, 0x4F))
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = GREEN
        p.font.name = "Consolas"
        p.space_after = Pt(2)
    return txBox

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def add_numbered_card(slide, left, top, width, height, number, title, desc, accent=ACCENT):
    card = add_shape_bg(slide, left, top, width, height, BG_CARD, RGBColor(0x30, 0x3A, 0x4F))
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), top + Inches(0.3), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(number)
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BG_DARK
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Calibri"
    add_text_box(slide, left + Inches(1), top + Inches(0.25), width - Inches(1.3), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(1), top + Inches(0.65), width - Inches(1.3), height - Inches(0.9),
                 desc, font_size=12, color=MID_GRAY)
    return card


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1),
             "Quintet", font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(3.3), Inches(2.3))

add_text_box(slide, Inches(2), Inches(3.6), Inches(9.3), Inches(0.6),
             "Unified Learning Platform", font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(0.5),
             "A Comprehensive Course Management System", font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.7), Inches(9.3), Inches(0.5),
             "Database Management Systems Lab Project", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

card = add_shape_bg(slide, Inches(3), Inches(5.5), Inches(7.3), Inches(1.5), BG_CARD, RGBColor(0x30, 0x3A, 0x4F))

add_text_box(slide, Inches(3.2), Inches(5.55), Inches(7), Inches(0.35),
             "Team Quintet", font_size=14, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3.2), Inches(5.9), Inches(7), Inches(0.4),
             "Sayon Sujit Mondal    |    Shresth M. Jha    |    Asmit Pandey", font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.2), Inches(6.3), Inches(7), Inches(0.4),
             "Harshvardhan Repaswal    |    Jay Jani", font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "Project Overview", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.8),
             "Quintet is a full-stack web application that provides a unified platform for course management, "
             "student enrollment, content delivery, and academic analytics -- built with modern technologies and "
             "a robust relational database design.",
             font_size=16, color=LIGHT_GRAY)

features = [
    ("Students", "Browse courses, enroll/unenroll,\nview materials, track grades", ACCENT),
    ("Instructors", "View assigned courses, manage\ncontent, grade enrolled students", GREEN),
    ("Analysts", "View database statistics, course\nsummaries, enrollment trends", PURPLE),
    ("Admins", "Manage all users, courses,\nenrollments, and assignments", ORANGE),
]

for i, (title, desc, color) in enumerate(features):
    left = Inches(0.8 + i * 3.1)
    top = Inches(2.6)
    w = Inches(2.8)
    h = Inches(2.2)
    card = add_shape_bg(slide, left, top, w, h, BG_CARD, color)
    add_text_box(slide, left, top + Inches(0.3), w, Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.9), w - Inches(0.4), Inches(0.9),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

stats = [("11", "Database Tables"), ("27+", "API Endpoints"), ("4", "User Roles"), ("10", "Sample Courses")]
for i, (num, label) in enumerate(stats):
    left = Inches(0.8 + i * 3.1)
    add_shape_bg(slide, left, Inches(5.2), Inches(2.8), Inches(1), BG_CARD)
    add_text_box(slide, left, Inches(5.25), Inches(2.8), Inches(0.5),
                 num, font_size=28, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(5.7), Inches(2.8), Inches(0.4),
                 label, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "Technology Stack", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

tech = [
    ("Backend", ACCENT, [
        ("FastAPI", "High-performance async Python web framework"),
        ("SQLAlchemy", "ORM for database interactions"),
        ("Pydantic", "Request/response validation"),
        ("python-jose", "JWT authentication tokens"),
        ("bcrypt", "Secure password hashing"),
        ("Uvicorn", "ASGI server"),
    ]),
    ("Frontend", GREEN, [
        ("Next.js 15", "React framework with SSR"),
        ("TypeScript", "Type-safe JavaScript"),
        ("Tailwind CSS", "Utility-first CSS framework"),
        ("shadcn/ui", "Accessible UI component library"),
        ("Recharts", "Data visualization & charts"),
        ("next-themes", "Dark/Light mode support"),
    ]),
    ("Database", PURPLE, [
        ("PostgreSQL", "Relational database engine"),
        ("Supabase", "Cloud-hosted Postgres"),
        ("11 Tables", "Normalized relational schema"),
        ("Foreign Keys", "Referential integrity"),
        ("Indexes", "Optimized query performance"),
        ("Cloud Pooler", "IPv4 connection pooling"),
    ]),
]

for i, (section, color, items) in enumerate(tech):
    left = Inches(0.8 + i * 4.1)
    top = Inches(1.5)
    w = Inches(3.8)

    card = add_shape_bg(slide, left, top, w, Inches(5.3), BG_CARD, color)
    add_text_box(slide, left, top + Inches(0.15), w, Inches(0.5),
                 section, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    for j, (name, desc) in enumerate(items):
        y = top + Inches(0.75 + j * 0.75)
        add_text_box(slide, left + Inches(0.3), y, w - Inches(0.6), Inches(0.3),
                     name, font_size=14, color=WHITE, bold=True)
        add_text_box(slide, left + Inches(0.3), y + Inches(0.25), w - Inches(0.6), Inches(0.3),
                     desc, font_size=11, color=MID_GRAY)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Database Schema -- Relational Design", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

tables = [
    ("USER", "user_id PK, email_id,\nrole, password", ACCENT, 0, 0),
    ("STUDENT", "student_id PK, user_id FK,\nage, skill_level,\ncategory, country", GREEN, 1, -1),
    ("INSTRUCTOR", "instructor_id PK, user_id FK,\nname, expertise", ORANGE, 1, 1),
    ("UNIVERSITY", "university_id PK,\nname, country", PURPLE, 2, 1.5),
    ("COURSE", "course_id PK, course_name,\nduration, program_type,\ninstructor_id FK,\nuniversity_id FK", ACCENT2, 2, 0),
    ("ENROLLMENT", "student_id FK,\ncourse_id FK,\nevaluation_score", GREEN, 2, -1.5),
    ("CONTENT", "content_id PK, course_id FK,\ntype, content_url", RED, 3, -0.5),
    ("TOPIC", "topic_id PK, topic_name", PURPLE, 3, 0.8),
    ("TEXTBOOK", "textbook_id PK, title,\nauthor, link", ORANGE, 3, -1.8),
]

base_x = Inches(1.5)
base_y = Inches(1.6)
x_step = Inches(3.0)

for tbl_name, attrs, color, col, row_offset in tables:
    x = base_x + col * x_step
    y = base_y + Inches(2.2) + row_offset * Inches(0.9)
    w = Inches(2.6)
    h = Inches(1.3)

    card = add_shape_bg(slide, x, y, w, h, BG_CARD, color)
    add_text_box(slide, x, y + Inches(0.05), w, Inches(0.3),
                 tbl_name, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.35), w - Inches(0.2), h - Inches(0.4),
                 attrs, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7), BG_CARD)
add_text_box(slide, Inches(1), Inches(6.35), Inches(11.3), Inches(0.6),
             "Junction Tables:   COURSE_TOPIC (course_id, topic_id)   |   TEXTBOOK_USED (course_id, textbook_id)   |   ENROLLMENT (student_id, course_id)       Total: 11 Tables",
             font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Backend Architecture", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

layers = [
    ("Client Request", "HTTP/JSON", ACCENT, Inches(0.5)),
    ("Routers", "URL to Handler mapping\n(auth, admin, students,\ninstructors, analyst,\ncourses, content)", GREEN, Inches(1.9)),
    ("Schemas", "Pydantic validation\nRequest/Response models\nType safety", ORANGE, Inches(3.3)),
    ("Services", "Business logic\nQuery construction\nError handling", PURPLE, Inches(4.7)),
    ("Models", "SQLAlchemy ORM\nTable definitions\nRelationships", RED, Inches(6.1)),
]

for name, desc, color, top in layers:
    if top > Inches(0.5):
        arrow_shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2.6), top - Inches(0.25), Inches(0.3), Inches(0.25))
        arrow_shape.fill.solid()
        arrow_shape.fill.fore_color.rgb = MID_GRAY
        arrow_shape.line.fill.background()

    card = add_shape_bg(slide, Inches(0.8), top + Inches(1), Inches(4.8), Inches(1.1), BG_CARD, color)
    add_text_box(slide, Inches(1), top + Inches(1.05), Inches(1.8), Inches(1),
                 name, font_size=15, color=color, bold=True)
    add_text_box(slide, Inches(2.8), top + Inches(1.05), Inches(2.5), Inches(1),
                 desc, font_size=11, color=MID_GRAY)

add_text_box(slide, Inches(6.2), Inches(1.3), Inches(6.5), Inches(0.4),
             "Request Flow Example", font_size=18, color=WHITE, bold=True)

code = '''# Router: receives HTTP request
@router.post("/api/auth/student/login")
async def login(data: UserLogin, db = Depends(get_db)):
    return login_user(db, data.email_id, data.password, "student")

class UserLogin(BaseModel):
    email_id: EmailStr
    password: str

def login_user(db, email, password, role):
    user = db.query(User).filter(
        User.email_id == email,
        User.role == role
    ).first()
    if not user or not verify_password(password, user.password):
        raise HTTPException(401, "Invalid credentials")
    token = create_access_token({"sub": user.email_id})
    return {"access_token": token, "role": role}

class User(Base):
    __tablename__ = "users"
    user_id = Column(Integer, primary_key=True)
    email_id = Column(String, unique=True, nullable=False)
    role = Column(String, nullable=False)
    password = Column(String, nullable=False)'''

add_code_block(slide, Inches(6.2), Inches(1.8), Inches(6.5), Inches(5.4), code, font_size=10)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "API Endpoints -- 27+ Routes", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

endpoint_groups = [
    ("Authentication", ACCENT, [
        "POST /api/auth/student/signup",
        "POST /api/auth/student/login",
        "POST /api/auth/instructor/login",
        "POST /api/auth/analyst/login",
        "POST /api/auth/admin/login",
    ]),
    ("Student", GREEN, [
        "GET  /api/students/profile",
        "GET  /api/students/courses",
        "POST /api/students/enroll",
        "DELETE /api/students/unenroll",
        "GET  /api/students/my-courses",
    ]),
    ("Admin", ORANGE, [
        "POST  /api/admin/instructors",
        "DELETE /api/admin/instructors/{id}",
        "POST  /api/admin/courses",
        "DELETE /api/admin/courses/{id}",
        "PUT   /api/admin/courses/assign",
        "GET/DELETE students, enrollments",
    ]),
    ("Instructor / Analyst", PURPLE, [
        "GET  /api/instructors/profile",
        "GET  /api/instructors/my-courses",
        "POST /instructors/courses/{id}/content",
        "GET  /api/analyst/statistics",
        "GET  /api/analyst/courses/summary",
        "GET  /api/analyst/enrollments/summary",
    ]),
]

for i, (group, color, endpoints) in enumerate(endpoint_groups):
    left = Inches(0.5 + i * 3.2)
    top = Inches(1.5)
    w = Inches(3.0)

    card = add_shape_bg(slide, left, top, w, Inches(5.3), BG_CARD, color)
    add_text_box(slide, left, top + Inches(0.1), w, Inches(0.4),
                 group, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    for j, ep in enumerate(endpoints):
        y = top + Inches(0.6 + j * 0.72)
        method = ep.split()[0]
        path = " ".join(ep.split()[1:])

        method_color = {
            "POST": GREEN, "GET": ACCENT, "PUT": ORANGE,
            "DELETE": RED, "GET/DELETE": PURPLE
        }.get(method, MID_GRAY)

        add_text_box(slide, left + Inches(0.15), y, Inches(0.8), Inches(0.25),
                     method, font_size=10, color=method_color, bold=True)
        add_text_box(slide, left + Inches(0.15), y + Inches(0.22), w - Inches(0.3), Inches(0.3),
                     path, font_size=10, color=LIGHT_GRAY)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Frontend Architecture", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(4), Inches(0.4),
             "Page Structure", font_size=20, color=GREEN, bold=True)

pages = [
    ("/", "Landing page with role cards"),
    ("/student/signup", "Student registration form"),
    ("/student/login", "Student authentication"),
    ("/student/dashboard", "Profile, courses, search, enroll/unenroll"),
    ("/instructor/login", "Instructor authentication"),
    ("/instructor/dashboard", "Profile, assigned courses"),
    ("/admin/login", "Admin authentication"),
    ("/admin/dashboard", "Manage users, courses, enrollments"),
    ("/analyst/login", "Analyst authentication"),
    ("/analyst/dashboard", "Statistics, charts, data summaries"),
    ("/courses/[id]", "Course detail, content, textbooks, grading"),
]

for i, (path, desc) in enumerate(pages):
    y = Inches(1.9 + i * 0.42)
    add_text_box(slide, Inches(0.8), y, Inches(2), Inches(0.4),
                 path, font_size=11, color=ACCENT, bold=True)
    add_text_box(slide, Inches(2.9), y, Inches(3), Inches(0.4),
                 desc, font_size=11, color=MID_GRAY)

add_text_box(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.4),
             "Key Features", font_size=20, color=ORANGE, bold=True)

features_list = [
    ("JWT Authentication", "Token-based auth with role-based access control"),
    ("Dark/Light Theme", "Persistent theme toggle with next-themes"),
    ("Search & Filter", "Real-time prefix search on course names"),
    ("Data Visualization", "Recharts for analytics dashboards"),
    ("Video Embedding", "YouTube video player embedded in course pages"),
    ("Textbook Links", "Direct links to reference materials"),
    ("Inline Grading", "Instructors grade students without page reload"),
    ("Content Management", "Instructors add/delete course materials"),
    ("Responsive Design", "Works on desktop and mobile screens"),
]

for i, (title, desc) in enumerate(features_list):
    y = Inches(1.9 + i * 0.55)
    card = add_shape_bg(slide, Inches(7), y, Inches(5.8), Inches(0.5), BG_CARD)
    add_text_box(slide, Inches(7.2), y + Inches(0.02), Inches(5.4), Inches(0.25),
                 title, font_size=12, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.2), y + Inches(0.24), Inches(5.4), Inches(0.22),
                 desc, font_size=10, color=MID_GRAY)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Role-Based Access Control", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5),
             "Each role has specific permissions enforced at the API level via JWT + dependency injection.",
             font_size=15, color=LIGHT_GRAY)

roles = [
    ("Student", ACCENT, [
        "  Self-registration (signup)",
        "  Browse all courses",
        "  Enroll / Unenroll from courses",
        "  View enrolled course materials",
        "  View own profile & grades",
        "  Cannot grade or manage users",
    ]),
    ("Instructor", GREEN, [
        "  No self-registration",
        "  View assigned courses",
        "  View enrolled students",
        "  Grade students (set scores)",
        "  Add/delete course content",
        "  Cannot manage other users",
    ]),
    ("Analyst", PURPLE, [
        "  No self-registration",
        "  View general DB statistics",
        "  View per-course summaries",
        "  View enrollment analytics",
        "  See charts and visualizations",
        "  Cannot modify any data",
    ]),
    ("Admin", ORANGE, [
        "  Predefined account only",
        "  Create/remove instructors",
        "  Create/delete courses",
        "  Assign instructors to courses",
        "  Enroll/drop students",
        "  View all student & instructor details",
    ]),
]

for i, (role, color, perms) in enumerate(roles):
    left = Inches(0.5 + i * 3.2)
    top = Inches(2.1)
    w = Inches(3.0)
    h = Inches(5.0)

    card = add_shape_bg(slide, left, top, w, h, BG_CARD, color)
    add_text_box(slide, left, top + Inches(0.1), w, Inches(0.5),
                 role, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    for j, perm in enumerate(perms):
        y = top + Inches(0.7 + j * 0.65)
        add_text_box(slide, left + Inches(0.2), y, w - Inches(0.4), Inches(0.55),
                     perm, font_size=12, color=LIGHT_GRAY)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Security & Authentication Flow", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

steps = [
    ("1", "Client sends credentials", "POST /api/auth/student/login\n{email_id, password}", ACCENT),
    ("2", "Server validates", "Check email exists, Verify role matches\nbcrypt.checkpw(password, hash)", GREEN),
    ("3", "JWT token created", "create_access_token({sub: email,\nrole: student, exp: 30min})", ORANGE),
    ("4", "Token returned to client", "Response: {access_token: 'eyJ...',\nrole: 'student', user_id: 6}", PURPLE),
    ("5", "Client stores token", "localStorage.setItem('token', jwt)\nAttach to all future requests", ACCENT2),
    ("6", "Protected endpoints", "Authorization: Bearer eyJ...\nget_current_user() verifies token", RED),
]

for i, (num, title, desc, color) in enumerate(steps):
    y = Inches(1.4 + i * 0.95)
    add_numbered_card(slide, Inches(0.8), y, Inches(6), Inches(0.85), num, title, desc, color)

add_text_box(slide, Inches(7.5), Inches(1.4), Inches(5), Inches(0.4),
             "Security Features", font_size=20, color=ORANGE, bold=True)

security_items = [
    ("bcrypt Hashing", "Passwords are never stored in plain text.\nbcrypt with salt rounds ensures security."),
    ("JWT Tokens", "Stateless authentication. Tokens expire\nafter 30 minutes for session safety."),
    ("Role Enforcement", "Depends(require_admin) blocks non-admin\nusers at the API layer automatically."),
    ("CORS Protection", "Only whitelisted origins (localhost:3000)\ncan make API requests."),
    ("SQL Injection Safe", "SQLAlchemy ORM with parameterized\nqueries prevents injection attacks."),
    ("Input Validation", "Pydantic schemas reject malformed\nrequests before hitting the database."),
]

for i, (title, desc) in enumerate(security_items):
    y = Inches(1.9 + i * 0.88)
    card = add_shape_bg(slide, Inches(7.5), y, Inches(5.3), Inches(0.8), BG_CARD, ACCENT)
    add_text_box(slide, Inches(7.7), y + Inches(0.05), Inches(4.9), Inches(0.3),
                 title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.7), y + Inches(0.32), Inches(4.9), Inches(0.45),
                 desc, font_size=10, color=MID_GRAY)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "SQL Queries via SQLAlchemy ORM", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.4),
             "Python ORM Code", font_size=18, color=GREEN, bold=True)

orm_code = '''# JOIN: Get student profile with user email
student = db.query(Student).filter(
    Student.user_id == user_id
).first()
user = db.query(User).filter(
    User.user_id == student.user_id
).first()

total_students = db.query(func.count(Student.student_id)).scalar()
avg_score = db.query(func.avg(Enrollment.evaluation_score)).scalar()

courses_per_uni = db.query(
    University.name,
    func.count(Course.course_id)
).join(Course, Course.university_id == University.university_id
).group_by(University.name).all()

enrollments = db.query(Enrollment, Student, User).join(
    Student, Enrollment.student_id == Student.student_id
).join(User, Student.user_id == User.user_id
).filter(Enrollment.course_id == course_id).all()'''

add_code_block(slide, Inches(0.8), Inches(1.9), Inches(5.8), Inches(5.0), orm_code, font_size=10)

add_text_box(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.4),
             "Equivalent SQL Generated", font_size=18, color=ORANGE, bold=True)

sql_code = '''-- JOIN: Student profile
SELECT s.*, u.email_id FROM students s
JOIN users u ON s.user_id = u.user_id
WHERE s.user_id = :user_id;

-- AGGREGATE: Statistics
SELECT COUNT(student_id) FROM students;
SELECT AVG(evaluation_score) FROM enrollments;

-- JOIN + GROUP BY: Courses per university
SELECT u.name, COUNT(c.course_id)
FROM universities u
JOIN courses c ON c.university_id = u.university_id
GROUP BY u.name;

-- MULTI-TABLE JOIN: Enrolled students
SELECT e.*, s.*, u.email_id
FROM enrollments e
JOIN students s ON e.student_id = s.student_id
JOIN users u ON s.user_id = u.user_id
WHERE e.course_id = :course_id;'''

add_code_block(slide, Inches(7), Inches(1.9), Inches(5.8), Inches(5.0), sql_code, font_size=10)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Analytics Dashboard -- Analyst View", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

stat_items = [
    ("Total Students", "5+", ACCENT),
    ("Total Courses", "10", GREEN),
    ("Total Instructors", "5", ORANGE),
    ("Avg Score", "72.3%", PURPLE),
    ("Universities", "5", RED),
]
for i, (label, val, color) in enumerate(stat_items):
    left = Inches(0.5 + i * 2.5)
    card = add_shape_bg(slide, left, Inches(1.4), Inches(2.3), Inches(1), BG_CARD, color)
    add_text_box(slide, left, Inches(1.45), Inches(2.3), Inches(0.4),
                 val, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(1.85), Inches(2.3), Inches(0.4),
                 label, font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(2.7), Inches(4), Inches(0.4),
             "Available Analytics", font_size=20, color=WHITE, bold=True)

analytics = [
    ("General Statistics", "Total students, courses, instructors, universities, enrollments; average evaluation score across all enrollments."),
    ("Course Summary", "Per-course breakdown: enrollment count, average score, instructor name, university -- displayed as bar charts via Recharts."),
    ("Enrollment Summary", "Total enrollments, average/min/max scores, top 5 most enrolled courses -- identifying popular courses and performance trends."),
    ("Geographic Distribution", "Students per country -- understanding the global reach of the learning platform."),
    ("University Analytics", "Courses per university -- comparing institutional offerings and course density."),
]

for i, (title, desc) in enumerate(analytics):
    y = Inches(3.2 + i * 0.8)
    card = add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(0.7), BG_CARD)
    add_text_box(slide, Inches(1), y + Inches(0.05), Inches(3), Inches(0.3),
                 title, font_size=14, color=ACCENT, bold=True)
    add_text_box(slide, Inches(1), y + Inches(0.33), Inches(11), Inches(0.35),
                 desc, font_size=11, color=MID_GRAY)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Demo -- How to Run", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.4),
             "Terminal 1: Backend (FastAPI)", font_size=18, color=GREEN, bold=True)

backend_cmd = '''cd backend
source .venv/bin/activate
uvicorn app.main:app --reload --port 8000


add_code_block(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.0), backend_cmd, font_size=12)

add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(0.4),
             "Terminal 2: Frontend (Next.js)", font_size=18, color=ACCENT, bold=True)

frontend_cmd = '''cd frontend
pnpm install
pnpm dev


add_code_block(slide, Inches(0.8), Inches(4.7), Inches(5.5), Inches(1.6), frontend_cmd, font_size=12)

add_text_box(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.4),
             "Test Credentials", font_size=18, color=ORANGE, bold=True)

creds = [
    ("Admin", "admin@quintet.com", "admin123", ORANGE),
    ("Student", "rahul.sharma@student.com", "student123", ACCENT),
    ("Instructor", "andrew.ng@quintet.com", "instructor123", GREEN),
    ("Analyst", "analyst1@quintet.com", "analyst123", PURPLE),
]

for i, (role, email, pwd, color) in enumerate(creds):
    y = Inches(1.9 + i * 1.1)
    card = add_shape_bg(slide, Inches(7), y, Inches(5.5), Inches(0.95), BG_CARD, color)
    add_text_box(slide, Inches(7.2), y + Inches(0.05), Inches(2), Inches(0.3),
                 role, font_size=15, color=color, bold=True)
    add_text_box(slide, Inches(7.2), y + Inches(0.32), Inches(5), Inches(0.3),
                 "Email:  " + email, font_size=12, color=LIGHT_GRAY)
    add_text_box(slide, Inches(7.2), y + Inches(0.58), Inches(5), Inches(0.3),
                 "Pass:   " + pwd, font_size=12, color=MID_GRAY)

add_shape_bg(slide, Inches(7), Inches(6.4), Inches(5.5), Inches(0.7), BG_CARD)
add_text_box(slide, Inches(7.2), Inches(6.45), Inches(5.1), Inches(0.6),
             "Run  python seed.py  to populate all tables with\n5 instructors, 5 students, 10 courses, 28 content items",
             font_size=12, color=MID_GRAY)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Challenges & Learnings", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.4),
             "Challenges Faced", font_size=20, color=RED, bold=True)

challenges = [
    ("IPv6 Connectivity", "Supabase direct connection uses IPv6 which was unavailable.\nSolution: Used Session Pooler (IPv4) connection string."),
    ("bcrypt + passlib Conflict", "passlib library incompatible with newer bcrypt versions.\nSolution: Replaced passlib with direct bcrypt usage."),
    ("CORS Configuration", "Frontend on different port blocked by browser CORS policy.\nSolution: Added all frontend origins to FastAPI CORS middleware."),
    ("Cascading Deletes", "Deleting a course/student required cleaning related tables first.\nSolution: Manual cascade delete order in service layer."),
]

for i, (title, desc) in enumerate(challenges):
    y = Inches(2.0 + i * 1.2)
    card = add_shape_bg(slide, Inches(0.8), y, Inches(5.5), Inches(1.05), BG_CARD, RED)
    add_text_box(slide, Inches(1), y + Inches(0.05), Inches(5.1), Inches(0.3),
                 title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(1), y + Inches(0.35), Inches(5.1), Inches(0.65),
                 desc, font_size=11, color=MID_GRAY)

add_text_box(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.4),
             "Key Learnings", font_size=20, color=GREEN, bold=True)

learnings = [
    ("Full-Stack Integration", "Connecting a React frontend to a Python backend\nvia REST APIs with JWT authentication."),
    ("Database Design", "Designing normalized relational schemas with proper\nforeign keys, constraints, and junction tables."),
    ("ORM vs Raw SQL", "Using SQLAlchemy ORM for safety while understanding\nthe SQL queries generated underneath."),
    ("Role-Based Security", "Implementing dependency injection for authentication\nand authorization at the API layer."),
]

for i, (title, desc) in enumerate(learnings):
    y = Inches(2.0 + i * 1.2)
    card = add_shape_bg(slide, Inches(7), y, Inches(5.5), Inches(1.05), BG_CARD, GREEN)
    add_text_box(slide, Inches(7.2), y + Inches(0.05), Inches(5.1), Inches(0.3),
                 title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.2), y + Inches(0.35), Inches(5.1), Inches(0.65),
                 desc, font_size=11, color=MID_GRAY)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "Future Scope", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

future_items = [
    ("Email Notifications", "Send enrollment confirmations and grade notifications via email.", ACCENT),
    ("Mobile Responsive", "Fully optimize all dashboards for mobile and tablet screens.", GREEN),
    ("Discussion Forums", "Add per-course discussion threads for student-instructor interaction.", ORANGE),
    ("Assignment Submission", "Allow instructors to create assignments and students to upload solutions.", PURPLE),
    ("AI Course Recommendations", "Recommend courses based on student skill level, category, and history.", ACCENT2),
    ("Advanced Analytics", "Add time-series analysis, completion rates, and predictive scoring models.", RED),
    ("Real-time Notifications", "WebSocket-based notifications for grades, enrollments, and announcements.", GREEN),
    ("Deployment", "Deploy backend to Railway/Render and frontend to Vercel for production.", ORANGE),
]

for i, (title, desc, color) in enumerate(future_items):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.3)
    top = Inches(1.5 + row * 1.35)
    w = Inches(5.8)
    h = Inches(1.15)

    card = add_shape_bg(slide, left, top, w, h, BG_CARD, color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), w - Inches(0.4), Inches(0.35),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.5), w - Inches(0.4), Inches(0.55),
                 desc, font_size=12, color=MID_GRAY)


slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1),
             "Thank You!", font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(3.8), Inches(2.3))

add_text_box(slide, Inches(2), Inches(4.1), Inches(9.3), Inches(0.6),
             "Quintet -- Unified Learning Platform", font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.7), Inches(9.3), Inches(0.5),
             "Database Management Systems Lab Project", font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

card = add_shape_bg(slide, Inches(3), Inches(5.5), Inches(7.3), Inches(1.5), BG_CARD, ACCENT)

add_text_box(slide, Inches(3.2), Inches(5.55), Inches(7), Inches(0.35),
             "Team Quintet", font_size=16, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3.2), Inches(6.0), Inches(7), Inches(0.35),
             "Sayon Sujit Mondal  |  Shresth M. Jha  |  Asmit Pandey", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.2), Inches(6.4), Inches(7), Inches(0.35),
             "Harshvardhan Repaswal  |  Jay Jani", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(6.9), Inches(9.3), Inches(0.3),
             "Backend: http://127.0.0.1:8000/docs    |    Frontend: http://localhost:3000",
             font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


output_path = "/home/shresth/Desktop/Qunitet_/presentation/Quintet_DBMS_Presentation.pptx"
prs.save(output_path)
print(f"\n  Presentation saved to: {output_path}")
print(f"  Total slides: {len(prs.slides)}")
print("\nSlide contents:")
slides_info = [
    "1.  Title Slide (Team Quintet)",
    "2.  Project Overview (4 roles + stats)",
    "3.  Technology Stack (Backend/Frontend/Database)",
    "4.  Database Schema (11 tables)",
    "5.  Backend Architecture (layered flow + code)",
    "6.  API Endpoints (27+ routes by category)",
    "7.  Frontend Architecture (pages + features)",
    "8.  Role-Based Access Control (permissions matrix)",
    "9.  Security & Authentication Flow (6-step JWT flow)",
    "10. SQL Queries via ORM (Python <-> SQL comparison)",
    "11. Analytics Dashboard (Analyst view breakdown)",
    "12. Demo -- How to Run (commands + credentials)",
    "13. Challenges & Learnings",
    "14. Future Scope (8 planned features)",
    "15. Thank You",
]
for s in slides_info:
    print(f"    {s}")
